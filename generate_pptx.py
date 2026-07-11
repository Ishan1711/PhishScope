from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Constants
BG_COLOR = RGBColor(10, 15, 30)       # Dark Navy/Blue Background
TITLE_COLOR = RGBColor(0, 255, 255)   # Cyan title
TEXT_COLOR = RGBColor(220, 220, 220)  # Light gray/white text
FONT_NAME = "Segoe UI"

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_notes(slide, notes):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes

def format_text_frame(text_frame, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = bold

def create_slide(title_text, body_text, notes, is_title=False):
    if is_title:
        slide_layout = prs.slide_layouts[0] # Title Slide
    else:
        slide_layout = prs.slide_layouts[1] # Title and Content
        
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    
    title = slide.shapes.title
    format_text_frame(title.text_frame, title_text, Pt(40), TITLE_COLOR, bold=True)
    
    if is_title:
        subtitle = slide.placeholders[1]
        format_text_frame(subtitle.text_frame, body_text, Pt(24), TEXT_COLOR)
    else:
        body = slide.placeholders[1]
        body.text_frame.clear()
        
        lines = body_text.strip().split('\n')
        for idx, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            p = body.text_frame.add_paragraph()
            # Determine indent level based on leading spaces/bullet formatting in my data
            if line.startswith('- '):
                p.text = line[2:]
                p.level = 0
            elif line.startswith('  - '):
                p.text = line[4:]
                p.level = 1
            else:
                p.text = line
                p.level = 0
                
            p.font.name = FONT_NAME
            p.font.size = Pt(20) if p.level == 0 else Pt(18)
            p.font.color.rgb = TEXT_COLOR

    if notes:
        add_notes(slide, notes)
    
    return slide

# SLIDE 1
create_slide(
    "PhishScope – AI Powered Email Header Analysis System",
    "Student Name: Ishan (or your name)\nUniversity: [Your University]\nDepartment: Cyber Security\nGuide: [Guide Name]\nSemester: Final Year",
    "Good morning everyone. I am presenting my final year project, PhishScope. It is an advanced, AI-powered system designed to automate and enhance the analysis of email headers to detect phishing threats.",
    is_title=True
)

# SLIDE 2
create_slide(
    "Problem Statement",
    "- Explosive Growth of Phishing: Phishing remains the primary vector for cyberattacks.\n- Email Spoofing Sophistication: Attackers manipulate headers to bypass traditional filters.\n- Complex Header Structures: Manual inspection is tedious, time-consuming, and prone to human error.\n- The Need for Automation: SOCs require intelligent tools to rapidly triage emails and extract Indicators of Compromise (IoCs).",
    "The core issue we address is the complexity of email spoofing. Attackers hide their tracks in raw email headers. Manually reading routing hops and authentication results is not scalable."
)

# SLIDE 3
create_slide(
    "Project Overview",
    "- What is PhishScope: A full-stack web application to automatically parse, analyze, and score email headers.\n- Core Objectives: Demystify complex headers, validate domain authentication, and leverage AI.\n- Main Modules:\n  - Input Module (Text & .eml files)\n  - Parsing Engine (Hop Extraction)\n  - AI Threat Analyzer (Groq)\n  - Reporting Module (PDF & JSON)\n- Expected Outcome: Comprehensive threat assessment dashboard.",
    "PhishScope bridges the gap between raw data and actionable intelligence. By uploading an email, the system extracts routing data and uses AI to explain the threat."
)

# SLIDE 4
create_slide(
    "System Architecture",
    "- User Interface: React & Vite-based frontend.\n- API Gateway: Flask backend routing secure HTTPS requests.\n- Processing Engines:\n  - Header Parser: Extracts Hops, IPs, and SPF/DKIM/DMARC.\n  - Threat Detection: Evaluates URLs and MITRE ATT&CK vectors.\n  - AI Analysis Engine: Groq API for natural language explanations.\n- Output Generation: Pydantic data structuring for PDFs via ReportLab.",
    "The frontend is built on React, which communicates with a Flask backend. The backend orchestrates local parsing logic and calls the external Groq AI API for advanced threat explanation."
)

# SLIDE 5
create_slide(
    "Key Features",
    "- Flexible Input Methods: Raw text or .eml/.txt file uploads.\n- Deep Header Parsing: Automated extraction of Hops, IPs, and headers.\n- Authentication Validation: Strict evaluation of SPF, DKIM, and DMARC.\n- URL Threat Assessment: Detection of shortened URLs and suspicious TLDs.\n- MITRE ATT&CK Mapping: Aligns detected threats with MITRE techniques.\n- AI-Powered Explanations: Translates anomalies into plain English.\n- Comprehensive Reporting: PDF and JSON investigation reports.",
    "PhishScope doesn't just read headers; it maps threats to MITRE ATT&CK frameworks, evaluates embedded URLs, and generates professional PDF reports for security teams."
)

# SLIDE 6
create_slide(
    "Technology Stack",
    "- Frontend: React 19, Vite, TypeScript, Framer Motion, Chart.js\n- Backend: Python 3, Flask, Waitress/Gunicorn\n- Data Modeling: Pydantic (Strict Data Validation)\n- AI & Machine Learning: Groq API (High-Speed LLM Inference)\n- Report Generation: ReportLab (PDF processing)\n- Deployment Platform: Render (Live Web Hosting)",
    "I utilized Vite and React for a responsive UI and Chart.js for visualization. The backend uses Python, Flask, Pydantic, and the Groq API for ultra-fast AI inference."
)

# SLIDE 7
create_slide(
    "Project Workflow",
    "- 1. Input & Validation: Backend sanitizes user input (.eml or text).\n- 2. Parsing: System extracts routing hops and metadata.\n- 3. Authentication Check: SPF, DKIM, DMARC evaluation.\n- 4. Threat Detection: URLs scanned, scoring algorithms calculate risk.\n- 5. AI Processing: Parsed data sent to Groq AI for explanation.\n- 6. Output Dashboard: Results and downloadable reports rendered via React.",
    "The workflow is seamless for the user. Behind the scenes, the header goes through validation, parsing, authentication checks, threat scoring, and AI processing."
)

# SLIDE 8
create_slide(
    "Advantages & Impact",
    "- Speed & Automation: Reduces analysis time from minutes to milliseconds.\n- High Accuracy: Combines static checks with AI contextual analysis.\n- Accessibility: Translates jargon into understandable insights for junior analysts.\n- Enterprise Ready: MITRE mapping and PDF exports for SOC integration.\n- Privacy Focused: Ephemeral processing; uploads are cleaned up immediately.",
    "The primary advantage is efficiency. It automates a highly manual task, ensuring accuracy while maintaining data privacy by discarding user uploads after processing."
)

# SLIDE 9
create_slide(
    "Future Scope",
    "- Browser Extensions: Direct integration into Outlook, Gmail for one-click analysis.\n- Real-Time Threat Intelligence: Integration with VirusTotal or AbuseIPDB.\n- Machine Learning Enhancement: Custom local ML model trained on phishing datasets.\n- Enterprise Dashboard: Multi-user support with historical tracking.\n- Automated Inbox Scanning: Background monitoring via IMAP integrations.",
    "The roadmap includes building browser extensions and integrating live threat intelligence APIs like VirusTotal to cross-reference extracted IPs and URLs."
)

# SLIDE 10
create_slide(
    "Conclusion",
    "- Problem Solved: Eliminated the complexity of manual email header analysis.\n- Objectives Achieved: Fast, automated, and AI-driven threat assessment tool.\n- Practical Impact: Empowers users to rapidly identify email spoofing.\n\nTHANK YOU & Questions?\n- GitHub: github.com/Ishan1711/PhishScope\n- Live Deployment: phishscope-zlzv.onrender.com",
    "In conclusion, PhishScope successfully automates complex header analysis, providing actionable intelligence. Thank you for your time. I am now open to any questions."
)

prs.save("PhishScope_Presentation.pptx")
print("Presentation generated successfully at PhishScope_Presentation.pptx")
